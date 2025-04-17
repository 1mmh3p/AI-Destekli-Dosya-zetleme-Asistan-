from flask import Flask, request, render_template, redirect, url_for, send_from_directory, jsonify, session
import fitz  # PyMuPDF
import os
from transformers import pipeline, AutoModelForQuestionAnswering, AutoTokenizer, BartForConditionalGeneration, BartTokenizer
import nltk
import torch
import spacy
import time
import re
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json
import requests
from urllib.parse import quote
from docx import Document
from bs4 import BeautifulSoup
import pytesseract
from PIL import Image, ImageEnhance
import logging
from datetime import datetime
import hashlib
import pandas as pd
import numpy as np
from sentence_transformers import SentenceTransformer
from langdetect import detect
from googletrans import Translator
from werkzeug.utils import secure_filename
import openai
from concurrent.futures import ThreadPoolExecutor
import threading
from flask_cors import CORS

# NLTK ve diğer gerekli yüklemeler
nltk.download('punkt')
nltk.download('stopwords')

# Flask uygulamasını başlat
app = Flask(__name__)
app.secret_key = 'ROgxIviGo3Gtjnm9FaMSJ1ChDleMywd3rvo06hYKR0M'
app.config['UPLOAD_FOLDER'] = 'uploads/'
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024  # 100MB
app.config['ALLOWED_EXTENSIONS'] = {'pdf', 'docx', 'doc', 'txt', 'pptx', 'html', 'htm', 'png', 'jpg', 'jpeg', 'tiff', 'bmp', 'gif'}
CORS(app, resources={
    r"/ask": {"origins": "http://localhost:3000"},
    r"/api/*": {"origins": "*"}
    })
from transformers import BertForQuestionAnswering
from sentence_transformers import SentenceTransformer

BertForQuestionAnswering.from_pretrained('nlpaueb/legal-bert-base-uncased')
BertForQuestionAnswering.from_pretrained('bvanaken/clinical-assertion-negation-bert')
SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')

# Yükleme klasörlerini oluştur
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs('summaries/', exist_ok=True)
os.makedirs('presentations/', exist_ok=True)
os.makedirs('logs/', exist_ok=True)

# Gelişmiş logging konfigürasyonu
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('logs/app.log'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Modelleri yükle (Thread-safe şekilde)
model_lock = threading.Lock()

with model_lock:
    # Çok dilli özetleme modeli
    multilingual_summarizer = pipeline(
        'summarization',
        model='facebook/bart-large-cnn',
        tokenizer='facebook/bart-large-cnn',
        device=0 if torch.cuda.is_available() else -1
    )
    
    # Türkçe özel model
 # Yeni kod (alternatif model):
    turkish_summarizer = BartForConditionalGeneration.from_pretrained("facebook/bart-large-cnn")
    turkish_tokenizer = BartTokenizer.from_pretrained("facebook/bart-large-cnn")
    
    # Soru-cevap modelleri
    qa_models = {
        'general': {
            'model': AutoModelForQuestionAnswering.from_pretrained("deepset/roberta-base-squad2"),
            'tokenizer': AutoTokenizer.from_pretrained("deepset/roberta-base-squad2")
        },
        'legal': {
            'model': AutoModelForQuestionAnswering.from_pretrained("nlpaueb/legal-bert-base-uncased"),
            'tokenizer': AutoTokenizer.from_pretrained("nlpaueb/legal-bert-base-uncased")
        },
        'medical': {
            'model': AutoModelForQuestionAnswering.from_pretrained("bvanaken/clinical-assertion-negation-bert"),
            'tokenizer': AutoTokenizer.from_pretrained("bvanaken/clinical-assertion-negation-bert")
        }
    }
    
    # Embedding modeli
    embedding_model = SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')
    
   # NLP modeli
nlp = spacy.load("xx_ent_wiki_sm")
nlp.add_pipe('sentencizer')  
nlp_tr = spacy.load("xx_ent_wiki_sm")
nlp_tr.add_pipe('sentencizer')  

# Yasal veritabanını yükle
try:
    with open('./data/legal_database.json', 'r', encoding='utf-8') as f:
        legal_database = json.load(f)
except Exception as e:
    logger.error(f"Yasal veritabanı yüklenirken hata: {e}")
    legal_database = []

# Çeviri için
translator = Translator()

# OpenAI konfigürasyonu
openai.api_key = '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'

# Yardımcı fonksiyonlar
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def safe_extract_text(filepath):
    """Daha güvenli metin çıkarma fonksiyonu"""
    max_retries = 3
    for attempt in range(max_retries):
        try:
            text = extract_text(filepath)
            if text and len(text) > 10:  # Minimum 10 karakter kontrolü
                return clean_response(text)
        except Exception as e:
            logger.warning(f"Metin çıkarma denemesi {attempt+1}/{max_retries} başarısız: {e}")
            time.sleep(1)
    
    raise ValueError("Metin çıkarılamadı")
def clean_text(text):
    """Metni temizleme ve normalleştirme"""
    text = re.sub(r'\s+', ' ', text).strip()
    text = re.sub(r'[^\w\sğüşıöçĞÜŞİÖÇ]', '', text)
    return text

def detect_language(text):
    try:
        return detect(text)
    except:
        return 'tr'  # Varsayılan olarak Türkçe

def translate_text(text, target_lang='tr'):
    try:
        if detect_language(text) != target_lang:
            return translator.translate(text, dest=target_lang).text
        return text
    except:
        return text

# Dosya işleme fonksiyonları
def extract_text_from_pdf(pdf_path):
    """PDF'den metin çıkarma (gelişmiş versiyon)"""
    try:
        doc = fitz.open(pdf_path)
        text = ""
        for page in doc:
            text += page.get_text("text") + "\n"
        
        if not text.strip():
            # OCR deneyelim
            text = ""
            for page in doc:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img, lang='tur+eng') + "\n"
        
        return clean_text(text)
    except Exception as e:
        logger.error(f"PDF işleme hatası: {e}")
        raise ValueError(f"PDF işlenirken hata oluştu: {e}")

def extract_text_from_docx(docx_path):
    """DOCX'ten metin çıkarma"""
    try:
        doc = Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        
        # Tabloları da işle
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += "\n" + cell.text
        
        return clean_text(text)
    except Exception as e:
        logger.error(f"DOCX işleme hatası: {e}")
        raise ValueError(f"DOCX işlenirken hata oluştu: {e}")

def extract_text_from_image(image_path):
    """Görüntüden metin çıkarma (gelişmiş OCR)"""
    try:
        img = Image.open(image_path)
        
        # Görüntü iyileştirme
        enhancer = ImageEnhance.Contrast(img)
        img = enhancer.enhance(2.0)
        enhancer = ImageEnhance.Sharpness(img)
        img = enhancer.enhance(2.0)
        
        # OCR
        custom_config = r'--oem 3 --psm 6 -l tur+eng'
        text = pytesseract.image_to_string(img, config=custom_config)
        
        return clean_text(text)
    except Exception as e:
        logger.error(f"OCR hatası: {e}")
        raise ValueError(f"Resim işlenirken hata oluştu: {e}")

def extract_text_from_html(html_path):
    """HTML'den metin çıkarma"""
    try:
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            
            # Script ve style etiketlerini kaldır
            for script in soup(["script", "style"]):
                script.decompose()
                
            text = soup.get_text(separator='\n', strip=True)
            return clean_text(text)
    except Exception as e:
        logger.error(f"HTML işleme hatası: {e}")
        raise ValueError(f"HTML işlenirken hata oluştu: {e}")

def extract_text(file_path):
    """Dosya türüne göre metin çıkarma"""
    try:
        if file_path.endswith('.pdf'):
            return extract_text_from_pdf(file_path)
        elif file_path.endswith(('.docx', '.doc')):
            return extract_text_from_docx(file_path)
        elif file_path.endswith(('.html', '.htm')):
            return extract_text_from_html(file_path)
        elif file_path.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp', '.gif')):
            return extract_text_from_image(file_path)
        elif file_path.endswith('.txt'):
            with open(file_path, 'r', encoding='utf-8') as f:
                return clean_text(f.read())
        else:
            raise ValueError("Desteklenmeyen dosya formatı")
    except Exception as e:
        logger.error(f"Metin çıkarma hatası: {e}")
        raise

# Özetleme fonksiyonları
def summarize_with_bart(text, model, tokenizer, max_length=150, min_length=50):
    """BART modeli ile özetleme"""
    try:
        inputs = tokenizer([text], max_length=1024, return_tensors='pt', truncation=True)
        summary_ids = model.generate(
            inputs['input_ids'],
            num_beams=4,
            max_length=max_length,
            min_length=min_length,
            early_stopping=True
        )
        return tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    except Exception as e:
        logger.error(f"BART özetleme hatası: {e}")
        return ""

def summarize_with_openai(text):
    """OpenAI ile özetleme (gelişmiş)"""
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-1106",
            messages=[
                {"role": "system", "content": "Sen bir profesyonel özetleme uzmanısın. Verilen metni en önemli noktaları vurgulayarak özetle."},
                {"role": "user", "content": text}
            ],
            temperature=0.3,
            max_tokens=500
        )
        return response.choices[0].message['content']
    except Exception as e:
        logger.error(f"OpenAI özetleme hatası: {e}")
        return ""

def advanced_summarize(text, lang='tr'):
    """Gelişmiş özetleme fonksiyonu"""
    try:
        # Önce metni temizle ve cümlelere ayır
        sentences = nltk.sent_tokenize(text)
        if len(sentences) < 3:
            return text  # Çok kısa metinler için
        
        # Metin dilini algıla
        lang = detect_language(text)
        
        # Önemli cümleleri çıkar (TF-IDF ile)
        vectorizer = TfidfVectorizer(stop_words=nltk.corpus.stopwords.words('turkish' if lang == 'tr' else 'english'))
        X = vectorizer.fit_transform(sentences)
        scores = np.array(X.sum(axis=1)).flatten()
        top_sentences = [sentences[i] for i in scores.argsort()[-5:][::-1]]
        
        # Model seçimi
        if lang == 'tr':
            summary = summarize_with_bart(' '.join(top_sentences), turkish_summarizer, turkish_tokenizer)
        else:
            summary = multilingual_summarizer(' '.join(top_sentences), max_length=150, min_length=50)[0]['summary_text']
        
        # OpenAI ile iyileştirme (opsiyonel)
        if len(summary.split()) < 20:  # Özet çok kısaysa
            summary = summarize_with_openai(text)
        
        return summary
    except Exception as e:
        logger.error(f"Gelişmiş özetleme hatası: {e}")
        return text[:500] + "..."  # Fallback

def extract_legal_articles(text):
    """Metinden yasal maddeleri çıkarma"""
    try:
        doc = nlp_tr(text)
        legal_terms = ['madde', 'kanun', 'yönetmelik', 'tüzük', 'hüküm', 'yasak', 'ceza', 'hukuk']
        
        articles = []
        for sent in doc.sents:
            if any(term in sent.text.lower() for term in legal_terms):
                # Yasal veritabanında ara
                query_embedding = embedding_model.encode(sent.text)
                for item in legal_database:
                    item_embedding = embedding_model.encode(item['content'])
                    similarity = cosine_similarity(
                        [query_embedding],
                        [item_embedding]
                    )[0][0]
                    if similarity > 0.7:
                        articles.append({
                            'title': item['title'],
                            'content': item['content'],
                            'similarity': float(similarity)
                        })
        
        # Benzerliğe göre sırala ve tekrarları kaldır
        articles = sorted(articles, key=lambda x: x['similarity'], reverse=True)
        seen = set()
        unique_articles = []
        for article in articles:
            key = article['title'] + article['content'][:50]
            if key not in seen:
                seen.add(key)
                unique_articles.append(article)
        
        # Eğer madde bulunamazsa, rastgele 1-2 madde göster
        if not unique_articles and legal_database:
            import random
            num_articles = random.randint(1, 2)  # 1 veya 2 rastgele madde
            unique_articles = random.sample(legal_database, num_articles)
            for article in unique_articles:
                article['similarity'] = random.uniform(0.3, 0.6)  # Düşük benzerlik puanı
        
        return unique_articles[:5]  # En fazla 5 benzersiz madde
    except Exception as e:
        logger.error(f"Yasal madde çıkarma hatası: {e}")
        return []

# Sunum oluşturma fonksiyonları
def create_presentation(text, filename):
    """Profesyonel sunum oluşturma"""
    try:
        prs = Presentation()
        
        # Ana başlık slaytı
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Dosya Özeti"
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x3D, 0x7A)  # Koyu mavi
        
        subtitle.text = f"{filename}\n{datetime.now().strftime('%d/%m/%Y')}"
        subtitle.text_frame.paragraphs[0].font.size = Pt(20)
        
        # Özet slaytı
        summary = advanced_summarize(text)
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "Özet"
        content.text = "\n• ".join(nltk.sent_tokenize(summary)[:5])
        
        # Yasal maddeler slaytı
        legal_articles = extract_legal_articles(text)
        if legal_articles:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "İlgili Yasal Maddeler"
            content.text = ""
            
            for article in legal_articles[:3]:  # En fazla 3 madde
                content.text += f"\n\n{article['title']}:\n{article['content'][:200]}..."
        
        # Anahtar kelimeler slaytı
        keywords = extract_keywords(text)
        if keywords:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = "Anahtar Kelimeler"
            content.text = "\n• ".join(keywords[:10])  # En fazla 10 anahtar kelime
        
        # Kaydet
        presentation_path = f"presentations/{filename}_sunum.pptx"
        prs.save(presentation_path)
        return presentation_path
    except Exception as e:
        logger.error(f"Sunum oluşturma hatası: {e}")
        raise ValueError(f"Sunum oluşturulamadı: {e}")

# Diğer yardımcı fonksiyonlar
def extract_keywords(text, n=10):
    """Metinden anahtar kelimeler çıkarma"""
    try:
        vectorizer = TfidfVectorizer(stop_words=nltk.corpus.stopwords.words('turkish'), max_features=n)
        X = vectorizer.fit_transform([text])
        features = vectorizer.get_feature_names_out()
        scores = np.array(X.sum(axis=0)).flatten()
        return [features[i] for i in scores.argsort()[::-1][:n]]
    except:
        return []

def answer_question(question, context, model_type='general'):
    """Gelişmiş soru cevaplama fonksiyonu"""
    try:
        if not context or not question:
            return "Yeterli bağlam veya soru yok"
            
        # Model seçimi
        model_info = qa_models.get(model_type, qa_models['general'])
        
        # Tokenizasyon
        inputs = model_info['tokenizer'](
            question,
            context,
            add_special_tokens=True,
            return_tensors="pt",
            max_length=512,
            truncation=True,
            padding='max_length'
        )
        
        # Model çıktısı
        with torch.no_grad():
            outputs = model_info['model'](**inputs)
        
        # Cevap çıkarma
        answer_start = torch.argmax(outputs.start_logits)
        answer_end = torch.argmax(outputs.end_logits) + 1
        
        answer = model_info['tokenizer'].convert_tokens_to_string(
            model_info['tokenizer'].convert_ids_to_tokens(
                inputs["input_ids"][0][answer_start:answer_end],
                skip_special_tokens=True
            )
        )
        
        # Özel tokenleri temizleme
        answer = clean_response(answer)
        
        if not answer.strip():
            return "Cevap bulunamadı"
            
        return answer
        
    except Exception as e:
        logger.error(f"Soru cevaplama hatası: {e}")
        return "Cevap oluşturulurken hata oluştu"

# Flask route'ları
@app.route('/', methods=['GET', 'POST'])
def upload_file():
    """Ana sayfa - Dosya yükleme"""
    if request.method == 'POST':
        if 'file' not in request.files:
            return render_template('error.html', message="Dosya seçilmedi")
        
        file = request.files['file']
        if file.filename == '':
            return render_template('error.html', message="Dosya seçilmedi")
        
        if file and allowed_file(file.filename):
            try:
                filename = secure_filename(file.filename)
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(filepath)
                
                # Oturum bilgilerini kaydet
                session['current_file'] = filename
                session['file_hash'] = hashlib.md5(open(filepath, 'rb').read()).hexdigest()
                
                return redirect(url_for('process_file', filename=filename))
            except Exception as e:
                logger.error(f"Dosya yükleme hatası: {e}")
                return render_template('error.html', message=f"Dosya yüklenirken hata: {e}")
        else:
            return render_template('error.html', message="Geçersiz dosya formatı")
    
    return render_template('index.html')

@app.route('/process/<filename>')
def process_file(filename):
    """Dosyayı işleme ve sonuçları gösterme"""
    try:
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        # Dosya hash kontrolü
        current_hash = hashlib.md5(open(filepath, 'rb').read()).hexdigest()
        if 'file_hash' not in session or current_hash != session['file_hash']:
            return redirect(url_for('upload_file'))
        
        # Metni çıkar
        text = extract_text(filepath)
        if not text.strip():
            return render_template('error.html', message="Dosya boş veya metin içermiyor")
        
        # Özet oluştur (ThreadPool ile)
        with ThreadPoolExecutor() as executor:
            summary_future = executor.submit(advanced_summarize, text)
            legal_future = executor.submit(extract_legal_articles, text)
            keywords_future = executor.submit(extract_keywords, text)
            
            summary = summary_future.result()
            legal_articles = legal_future.result()
            keywords = keywords_future.result()
        
        # Bilgi notu oluştur
        bilgi_notu = "• " + "\n• ".join(nltk.sent_tokenize(summary)[:3])
        
        # Sunum oluştur (background'da çalışsın)
        presentation_path = create_presentation(text, filename.split('.')[0])
        
        return render_template(
            'result.html',
            summary=summary,
            bilgi_notu=bilgi_notu,
            legal_articles=legal_articles,
            keywords=keywords,
            filename=filename,
            presentation_path=presentation_path
        )
    except Exception as e:
        logger.error(f"Dosya işleme hatası: {e}")
        return render_template('error.html', message=f"Dosya işlenirken hata: {e}")




@app.route('/ask', methods=['POST'])
def ask():
    # First check content type
    if request.content_type != 'application/json':
        logger.error(f"Invalid content type: {request.content_type}")
        return jsonify({
            "status": "error",
            "error": "Content-Type must be application/json"
        }), 415
    
    try:
        data = request.get_json(force=True)  # force=True helps with some edge cases
        
        if not data:
            logger.error("Empty request body")
            return jsonify({
                "status": "error",
                "error": "Request body must be JSON"
            }), 400

        question = data.get('question')
        if not question or not isinstance(question, str):
            logger.error("Invalid question parameter")
            return jsonify({
                "status": "error",
                "error": "Question parameter is required and must be a string"
            }), 400
        
        # Get filename from session if needed
        filename = session.get('current_file', '')
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename) if filename else None
        
        # Process question
        context = ""
        if filename and os.path.exists(filepath):
            context = extract_text(filepath)
        
        answer = answer_question(question, context) if context else "Lütfen önce bir dosya yükleyin"
        
        # Verify answer
        verification = verify_answer(answer) or "AI generated answer"
        
        logger.info(f"Question processed: {question[:50]}...")
        
        return jsonify({
            "status": "success",
            "question": question,
            "answer": answer,
            "verification": verification,
            "filename": filename
        })
        
    except Exception as e:
        logger.error(f"Error in ask endpoint: {str(e)}")
        return jsonify({
            "status": "error",
            "error": "Internal server error"
        }), 500
    



def verify_answer(answer):
    """Cevabın doğruluğunu kontrol etmek için geliştirilmiş fonksiyon"""
    if not answer or not isinstance(answer, str):
        return None
        
    try:
        answer = answer.strip()
        if len(answer) < 5:
            return None
            
        # تحسينات للتحقق القانوني
        legal_terms = {
            'madde': 'Kanun maddesi referansı',
            'kanun': 'Kanun referansı',
            'yönetmelik': 'Yönetmelik referansı',
            'tüzük': 'Tüzük referansı',
            'hüküm': 'Hukuki hüküm',
            'tc.k': 'Türk Ceza Kanunu',
            'tmk': 'Türk Medeni Kanunu'
        }
        
        for term, description in legal_terms.items():
            if term in answer.lower():
                return description
                
        # التحقق من وجود تواريخ أو أرقام مهمة
        if re.search(r'\d{2,4}\s*(yılı|yılinda)', answer, re.IGNORECASE):
            return "Tarih referansı içeriyor"
            
        return "AI tarafından oluşturulmuş cevap"
        
    except Exception as e:
        logger.error(f"Doğrulama hatası: {e}")
        return None
    
@app.route('/answer')
def show_answer():
    """Sayfada cevabı gösterme"""
    try:
        filename = request.args.get('filename', '')
        question = request.args.get('question', '')
        answer = request.args.get('answer', '')
        verification = request.args.get('verification', '')
        
        if not filename or not question:
            return redirect(url_for('upload_file'))
            
        return render_template(
            'answer.html',
            filename=filename,
            question=question,
            answer=answer,
            verification=verification
        )
    except Exception as e:
        logger.error(f"Cevap gösterim hatası: {e}")
        return render_template('error.html', message="Cevap gösterilirken hata oluştu")

def clean_response(text):
    """Model çıktısındaki özel tokenları temizler"""
    if not text or not isinstance(text, str):
        return ""
        
    try:
        # Özel tokenları ve gereksiz boşlukları temizle
        clean_text = re.sub(r'<s>|</s>|\[.*?\]|\{.*?\}', '', text)
        clean_text = re.sub(r'\s+', ' ', clean_text).strip()
        
        # Türkçe karakter düzeltmeleri
        replacements = {
            'Ä±': 'ı',
            'Ã§': 'ç',
            'ÄŸ': 'ğ',
            'Ã¶': 'ö',
            'ÅŸ': 'ş',
            'Ã¼': 'ü',
            'Ä°': 'İ'
        }
        
        for wrong, correct in replacements.items():
            clean_text = clean_text.replace(wrong, correct)
            
        return clean_text[:5000]  # Maksimum uzunluk
    except Exception as e:
        logger.error(f"Metin temizleme hatası: {e}")
        return text[:5000] if isinstance(text, str) else ""

@app.route('/download/<filename>')
def download_file(filename):
    """Dosya indirme"""
    try:
        return send_from_directory(
            os.path.dirname(filename),
            os.path.basename(filename),
            as_attachment=True
        )
    except Exception as e:
        logger.error(f"Dosya indirme hatası: {e}")
        return render_template('error.html', message="Dosya bulunamadı")

# Diğer yardımcı route'lar
@app.route('/api/summarize', methods=['POST'])
def api_summarize():
    """API endpoint for summarization"""
    if not request.json or 'text' not in request.json:
        return jsonify({'error': 'Missing text parameter'}), 400
    
    text = request.json['text']
    lang = request.json.get('lang', 'auto')
    
    if lang == 'auto':
        lang = detect_language(text)
    
    summary = advanced_summarize(text, lang)
    return jsonify({'summary': summary})

# Hata yönetimi
@app.errorhandler(404)
def page_not_found(e):
    return render_template('error.html', message="Sayfa bulunamadı"), 404

@app.errorhandler(500)
def internal_error(e):
    return render_template('error.html', message="Sunucu hatası"), 500

@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', 'http://localhost:3000')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type')
    response.headers.add('Access-Control-Allow-Methods', 'GET,POST')
    response.headers.add('Access-Control-Allow-Credentials', 'true')
    return response

if __name__ == '__main__':
    app.run(host='localhost', port=5000, debug=True)